import pandas as pd
from deep_translator import GoogleTranslator
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from sklearn.feature_extraction.text import CountVectorizer
from sklearn.naive_bayes import MultinomialNB
from sklearn.pipeline import make_pipeline
import re

def load_dictionaries(csv_files):
    """
    Load translation dictionaries from uploaded CSV files.

    Args:
        csv_files: List of uploaded CSV files.

    Returns:
        A dictionary containing word/phrase translations.
    """
    translation_dict = {}
    for file in csv_files:
        df = pd.read_csv(file)
        for index, row in df.iterrows():
            for col1 in df.columns:
                for col2 in df.columns:
                    if col1 != col2:
                        source_word = row[col1]
                        target_word = row[col2]
                        translation_dict[source_word] = target_word
                        translation_dict[target_word] = source_word
    return translation_dict

def train_translation_model(csv_files):
    """
    Train translation models using uploaded CSV files.

    Args:
        csv_files: List of uploaded CSV files.

    Returns:
        A dictionary containing trained translation models.
    """
    data = []
    for file in csv_files:
        df = pd.read_csv(file)
        data.append(df)
    data = pd.concat(data, ignore_index=True)
    
    models = {}
    for col1 in data.columns:
        for col2 in data.columns:
            if col1 != col2:
                X = data[col1]
                y = data[col2]
                
                model = make_pipeline(CountVectorizer(), MultinomialNB())
                model.fit(X, y)
                
                models[(col1, col2)] = model
                
    return models

def translate_with_model(text, source_lang, target_lang, models):
    """
    Translate text using a trained model.

    Args:
        text: Text to translate.
        source_lang: Source language.
        target_lang: Target language.
        models: Dictionary of trained models.

    Returns:
        Translated text.
    """
    model_key = (source_lang, target_lang)
    if model_key in models:
        model = models[model_key]
        return model.predict([text])[0]
    else:
        return "Model only supports translations between the specified languages."

def translate_text(text, translation_dict):
    """
    Translate text using the provided translation dictionary.

    Args:
        text: Text to translate.
        translation_dict: Dictionary containing word/phrase translations.

    Returns:
        Translated text.
    """
    if not text:
        return text
    words = text.split()
    translated_words = [translation_dict.get(word, word) for word in words]
    return ' '.join(translated_words)

def translate_with_google(text, target_lang, translation_dict):
    """
    Translate text using Google Translate, with an initial pass through the translation dictionary.

    Args:
        text: Text to translate.
        target_lang: Target language.
        translation_dict: Dictionary containing word/phrase translations.

    Returns:
        Translated text.
    """
    text = translate_text(text, translation_dict)
    if not text.strip():
        return text
    translator = GoogleTranslator(source='auto', target=target_lang)
    translated = translator.translate(text)
    return translated

def translate_paragraph(paragraph, target_lang, translation_dict, models=None):
    """
    Translate a paragraph, updating its runs with translated text.

    Args:
        paragraph: Paragraph object to translate.
        target_lang: Target language.
        translation_dict: Dictionary containing word/phrase translations.
        models: Dictionary of trained models (optional).

    Returns:
        Translated paragraph.
    """
    for run in paragraph.runs:
        if models:
            translated_text = translate_with_model(run.text, 'auto', target_lang, models)
        else:
            translated_text = translate_with_google(run.text, target_lang, translation_dict)
        run.text = translated_text if translated_text is not None else ""
    return paragraph

def translate_word(doc, target_lang, translation_dict, models=None):
    """
    Translate a Word document.

    Args:
        doc: Document object to translate.
        target_lang: Target language.
        translation_dict: Dictionary containing word/phrase translations.
        models: Dictionary of trained models (optional).

    Returns:
        Translated document.
    """
    for paragraph in doc.paragraphs:
        translate_paragraph(paragraph, target_lang, translation_dict, models)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    translate_paragraph(paragraph, target_lang, translation_dict, models)
    return doc

def translate_excel(wb, target_lang, translation_dict, models=None):
    """
    Translate an Excel workbook.

    Args:
        wb: Workbook object to translate.
        target_lang: Target language.
        translation_dict: Dictionary containing word/phrase translations.
        models: Dictionary of trained models (optional).

    Returns:
        Translated workbook.
    """
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value:
                    if models:
                        cell.value = translate_with_model(cell.value, 'auto', target_lang, models)
                    else:
                        cell.value = translate_with_google(cell.value, target_lang, translation_dict)
    return wb

def translate_pptx(prs, target_lang, translation_dict, models=None):
    """
    Translate a PowerPoint presentation, maintaining the original text alongside the translation.

    Args:
        prs: Presentation object to translate.
        target_lang: Target language.
        translation_dict: Dictionary containing word/phrase translations.
        models: Dictionary of trained models (optional).

    Returns:
        Translated presentation.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if models:
                            translated_text = translate_with_model(run.text, 'auto', target_lang, models)
                        else:
                            translated_text = translate_with_google(run.text, target_lang, translation_dict)
                        if translated_text:
                            run.text = f"{run.text} / {translated_text}"
    return prs

def translate_srt(srt_content, target_lang, translation_dict, models=None):
    """
    Translate an SRT subtitle file.

    Args:
        srt_content: Content of the SRT file to translate.
        target_lang: Target language.
        translation_dict: Dictionary containing word/phrase translations.
        models: Dictionary of trained models (optional).

    Returns:
        Translated SRT content.
    """
    translated_lines = []
    for line in srt_content.splitlines():
        if re.match(r'\d{2}:\d{2}:\d{2}', line):
            translated_lines.append(line)
        else:
            if models:
                translated_text = translate_with_model(line, 'auto', target_lang, models)
            else:
                translated_text = translate_with_google(line, target_lang, translation_dict)
            translated_lines.append(translated_text if translated_text is not None else line)
    return '\n'.join(translated_lines)

def save_translated_file(translated_file, file_type, initial_name, target_lang):
    """
    Save the translated file.

    Args:
        translated_file: Translated file object.
        file_type: Type of the file (docx, xlsx, pptx, srt).
        initial_name: Initial name of the file.
        target_lang: Target language.

    Returns:
        The name of the saved file.
    """
    file_name = f"{initial_name}_{target_lang}.{file_type}"
    if file_type == "docx":
        translated_file.save(file_name)
    elif file_type == "xlsx":
        translated_file.save(file_name)
    elif file_type == "pptx":
        translated_file.save(file_name)
    elif file_type == "srt":
        with open(file_name, "w", encoding="utf-8") as f:
            f.write(translated_file)
    return file_name
